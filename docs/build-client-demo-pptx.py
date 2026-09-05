#!/usr/bin/env python3
"""Build the SARA ROSE client demonstration PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x0F, 0x14)
PANEL = RGBColor(0x14, 0x1A, 0x22)
GOLD = RGBColor(0xD4, 0xA0, 0x37)
WHITE = RGBColor(0xF5, 0xF0, 0xE6)
MUTED = RGBColor(0xB8, 0xB0, 0xA0)
LINE = RGBColor(0x2A, 0x32, 0x3C)

OUT = Path(__file__).resolve().parent / "SARA-ROSE-Client-Demo.pptx"


def set_run(run, text, size=18, bold=False, color=WHITE, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def paint_bg(slide, prs):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    foot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        prs.slide_height - Inches(0.42),
        prs.slide_width,
        Inches(0.42),
    )
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0x08, 0x0B, 0x0F)
    foot.line.fill.background()
    box = slide.shapes.add_textbox(
        Inches(0.4), prs.slide_height - Inches(0.38), Inches(12.4), Inches(0.32)
    )
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    set_run(p.add_run(), "SARA ROSE NIGERIA LIMITED  ·  Confidential client demonstration  ·  Sagamu, Ogun State", 10, False, MUTED)


def kicker(slide, text, top=0.28):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.4), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    set_run(p.add_run(), text.upper(), 12, True, GOLD)


def title(slide, text, top=0.52, size=32):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.4), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, size, True, WHITE)


def bullets(slide, items, left=0.5, top=1.35, width=12.4, height=5.4, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        set_run(p.add_run(), "▸  " + item, size, False, WHITE)


def card(slide, x, y, w, h, heading, body):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    hb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.36), Inches(0.4))
    p = hb.text_frame.paragraphs[0]
    set_run(p.add_run(), heading, 16, True, GOLD)
    bb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.52), Inches(w - 0.36), Inches(h - 0.7))
    tf = bb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), body, 14, False, WHITE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Client demonstration")
    title(s, "SARA ROSE Nigeria Limited", 0.85, 40)
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), "A live website for a heavy-equipment trader in Sagamu —\ncompany story, equipment portfolio, enquiry desk, and an administration panel.", 22, False, MUTED)
    card(s, 0.5, 3.3, 3.9, 2.4, "The business", "Trader in construction and industrial machinery since 2012. One industry. Direct dealing.")
    card(s, 4.7, 3.3, 3.9, 2.4, "The product", "Public website plus a signed-in dashboard the company can operate after the demo.")
    card(s, 8.9, 3.3, 3.9, 2.4, "The contact", "Mr. Akram Haider\n+234 80 6665 1111\ncontact@sararose.com")
    add_notes(s, "Open with the company name and the purpose of the meeting: this is the working website, not a brochure mock-up. Mention Sagamu / Ogun State and that the named contact remains Mr. Akram Haider.")

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "How we will spend the next ten minutes")
    title(s, "Agenda")
    bullets(s, [
        "Who SARA ROSE is, and why this site exists",
        "Walk the public pages a customer would use",
        "Show the five equipment groups and an enquiry",
        "Register / sign in, then open the dashboard",
        "Show what the company can change itself: home slider and header links",
        "Questions and next steps",
    ])
    add_notes(s, "Keep this slide short. Promise a live click-through rather than a long talk.")

    # 3 Company
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "The organisation")
    title(s, "A specialist trader, not a general merchant")
    bullets(s, [
        "SARA ROSE NIGERIA LIMITED has traded heavy equipment since 2012.",
        "Head office: Km 12, Sagamu–Benin Express Way, opposite Navy Merchant, Ogun State.",
        "Business type: Trader. Industry: heavy / construction and industrial equipment.",
        "The work starts with the requirement — ground, timeline, duty — then the machine category.",
        "Brands, models and availability are confirmed at enquiry, not from a stock list on the site.",
        "Every commercial conversation has a named person: Mr. Akram Haider.",
    ])
    add_notes(s, "This language matches the company profile. Stress accountability and specialisation.")

    # 4 Why the site
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Why we built this")
    title(s, "A digital front door that matches how SARA ROSE actually deals")
    bullets(s, [
        "Customers can read the company, the portfolio and the values before they call.",
        "An enquiry is written down in the database — it does not disappear in a private inbox.",
        "Registration and login identify returning customers and staff.",
        "The dashboard lets the company update the home images and the menu without waiting on a developer for every change.",
        "The public copy stays honest: this is a trader’s site, not a manufacturer catalogue.",
    ])
    add_notes(s, "Position the site as an operating tool, not decoration.")

    # 5 Solution map
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "What was delivered")
    title(s, "One application, two sides")
    card(s, 0.5, 1.4, 6.0, 4.8, "Public website",
         "Home slider and company story\nAbout, Why us, Vision & values\nFive equipment categories, eleven machine types\nEnquire form\nRegistration and login\nDynamic header links")
    card(s, 6.8, 1.4, 6.0, 4.8, "Administration panel",
         "Signed-in dashboard\nHome slider upload and order\nHeader link manager\nRegistration list\nEnquiry records via the API\nRoles: Admin, Staff, User")
    add_notes(s, "Do not linger. This is the map before you open the browser.")

    # 6 Public tour
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Live demo — part one")
    title(s, "Open the public site")
    bullets(s, [
        "Address: http://127.0.0.1:43123  (or the hosted URL you are using today)",
        "Home: slider, logo, Sagamu, buttons for portfolio, enquiry and login.",
        "About: trading since 2012 and how the company works.",
        "Why us: longevity, focus, portfolio, local office, named contact.",
        "Vision & values: reliability, integrity, customer focus, professionalism, quality, partnerships.",
        "Pause on the named contact so the room sees a person, not a generic form.",
    ])
    add_notes(s, "Click slowly. Read one sentence from About aloud. Do not skip Why us — it is the commercial argument.")

    # 7 Equipment
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Live demo — part two")
    title(s, "Equipment portfolio")
    card(s, 0.5, 1.4, 2.4, 2.15, "Earthmoving", "Excavators\nBulldozers\nWheel loaders")
    card(s, 3.1, 1.4, 2.4, 2.15, "Construction", "Backhoe loaders\nMotor graders")
    card(s, 5.7, 1.4, 2.4, 2.15, "Material handling", "Forklifts\nOther machinery")
    card(s, 8.3, 1.4, 2.4, 2.15, "Road & compaction", "Rollers\nSkid-steer rollers")
    card(s, 10.9, 1.4, 2.4, 2.15, "Transport & lifting", "Dump trucks\nCranes")
    bullets(s, [
        "Open Earthmoving, then Excavators. Say: specification is confirmed when they enquire.",
        "Use “Start an enquiry” from a machine page so the audience sees the path from machine to conversation.",
    ], top=3.8, height=2.4, size=16)
    add_notes(s, "Do not pretend you have a price list. The profile forbids it.")

    # 8 Enquiry
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Live demo — part three")
    title(s, "Capture an enquiry")
    bullets(s, [
        "Page: Enquire (/contact).",
        "Fill name, telephone, email, a machine type, and a short requirement (at least a sentence).",
        "Submit. The record is stored in MySQL for follow-up.",
        "Commercially, the conversation still belongs to Mr. Akram Haider.",
        "Later, from the dashboard or API, staff can see the same row — nothing is lost in a personal inbox.",
    ])
    add_notes(s, "Use a realistic example: excavator for a foundation pit in Sagamu. Keep it short.")

    # 9 Accounts
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Live demo — part four")
    title(s, "Registration and login")
    bullets(s, [
        "Registration asks for a username, name, email, telephone, user type, role and password.",
        "That creates a profile and a login row. Public users cannot make themselves Admin.",
        "Login checks username and password against the userMaster table (hash and normal password columns).",
        "A small number captcha sits on the login form.",
        "After a correct login, choose Continue to dashboard.",
    ])
    add_notes(s, "If time is short, skip creating a new user and go straight to the admin account on the next slide.")

    # 10 Admin login
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Demonstration account")
    title(s, "Sign in as the seeded administrator")
    card(s, 0.5, 1.45, 6.1, 4.6, "Use these details on /login",
         "Username:  admin\nPassword:  SaraRose_Admin_2024\n\nSolve the captcha (for example 2 + 3 = 5).\nClick Login.\nWhen asked, continue to the dashboard.\n\nRole: Admin   ·   User type: Internal")
    card(s, 6.9, 1.45, 5.9, 4.6, "What the room should see",
         "Sidebar shows username admin.\nGreeting uses the administrator’s name.\nTiles show categories, machines, enquiries and registrations.\nRecent registrations list includes the admin (and any user you just created).")
    add_notes(s, "Type the password slowly or paste it. If captcha fails, click New sum. Do not show NormalPassword from the database on screen.")

    # 11 Dashboard
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Live demo — part five")
    title(s, "The dashboard the company will use")
    bullets(s, [
        "Overview: the day’s greeting, role tags, company facts, quick actions.",
        "Home slider: add an image, view the set, change order or remove a slide. The public home page picks this up when dynamic mode is on.",
        "Header links: add, hide, reorder, or mark the Enquire button. The public header follows the table.",
        "Registrations: people who filled the public form — names, usernames, roles — never passwords.",
        "Sign out returns the browser to login.",
    ])
    add_notes(s, "If you have a sample site photo, upload it as a slider image. That is the most convincing admin moment.")

    # 12 Technology
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "For the technical stakeholder in the room")
    title(s, "How it is put together")
    card(s, 0.5, 1.4, 4.0, 4.6, "Website", "Angular 19\nRuns on port 43123\nPublic pages and the signed-in shell")
    card(s, 4.7, 1.4, 4.0, 4.6, "API", "ASP.NET Core 8\nPort 43124\nSwagger at /swagger\nJSON for company, catalogue, login")
    card(s, 8.9, 1.4, 4.0, 4.6, "Database", "MySQL 8 — database sararose\nCatalogue, enquiries, slider, header, registrations, userMaster")
    add_notes(s, "Only if someone asks. Otherwise skip. Mention Workbench script backend/sql/update-all-db.sql for a full rebuild.")

    # 13 Data / login table
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Login data")
    title(s, "userMaster is the table that decides who gets in")
    bullets(s, [
        "Username (unique), email, name, phone, role, user type.",
        "HashPassword (SHA-256) and NormalPassword — both checked at login.",
        "Active flag. Created date. No RegistrationId column.",
        "The login API never returns the normal password to the browser.",
        "Company profile paragraphs are not in SQL; they come from the API as in the client document.",
    ])
    add_notes(s, "Useful if a director asks where passwords live. Keep it one minute.")

    # 14 Click path cheat sheet
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Presenter cheat sheet")
    title(s, "Click path if the network is slow")
    bullets(s, [
        "Home → Equipment → Excavators → Enquire (submit one form).",
        "Why us (one scroll) → Vision & values (one scroll).",
        "Login as admin → Continue to dashboard.",
        "Dashboard → Add slide (optional) → Header links (point at Enquire) → Registrations.",
        "Sign out. That is a complete demonstration.",
        "If the API is down, show Swagger health only after the room has seen the website.",
    ], size=17)
    add_notes(s, "Print this slide or keep it on a second screen. It is the recovery plan.")

    # 15 Close
    s = prs.slides.add_slide(blank)
    paint_bg(s, prs)
    kicker(s, "Close")
    title(s, "Ready for the work ahead", 1.6, 36)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.2), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), "SARA ROSE NIGERIA LIMITED now has a working site: the company story,\nthe five equipment groups, an enquiry desk, and a dashboard the team can run.", 20, False, WHITE)
    p = tf.add_paragraph()
    p.space_before = Pt(16)
    set_run(p.add_run(), "Questions. Then we agree hosting, content photos, and who will operate the slider.", 18, False, MUTED)
    box = s.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(12.2), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Mr. Akram Haider  ·  +234 80 6665 1111  ·  contact@sararose.com", 16, True, GOLD)
    p = tf.add_paragraph()
    set_run(p.add_run(), "Km 12, Sagamu–Benin Express Way, Ogun State, Nigeria", 14, False, MUTED)
    add_notes(s, "Stop talking. Take questions. Offer the documentation PDF/Markdown in docs/.")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
